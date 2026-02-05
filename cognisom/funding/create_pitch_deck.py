#!/usr/bin/env python3
"""
Generate eyentelligence pitch deck with deep-tech styling (B2)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color palette
COLORS = {
    'primary_dark': RGBColor(15, 37, 56),      # #0F2538
    'gradient_blue': RGBColor(31, 111, 169),   # #1F6FA9
    'gradient_teal': RGBColor(42, 166, 201),   # #2AA6C9
    'text_white': RGBColor(255, 255, 255),     # #FFFFFF
    'text_light': RGBColor(200, 208, 221),     # #C8D0DD
    'text_accent': RGBColor(136, 169, 209),    # #88A9D1
    'bio_accent': RGBColor(120, 194, 164),     # #78C2A4
    'background': RGBColor(26, 29, 37),        # #1A1D25
}

def create_gradient_background(slide, color1, color2):
    """Add gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 90.0
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def add_title_text(slide, text, top, left, width, height, font_size, color, bold=True):
    """Add title text box"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = 'Arial'  # Fallback for Poppins
    
    return textbox

def add_body_text(slide, text, top, left, width, height, font_size, color, align=PP_ALIGN.LEFT):
    """Add body text box"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.name = 'Arial'
        paragraph.space_after = Pt(12)
    
    return textbox

def add_bullet_text(slide, items, top, left, width, height, font_size, color):
    """Add bulleted text"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(8)
    
    return textbox

def create_presentation():
    """Create the full pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'eyentelligence', 
                   Inches(1.5), Inches(1), Inches(8), Inches(1), 
                   72, COLORS['text_white'])
    add_title_text(slide, 'Biological Simulation + Humanoid Cognitive Systems',
                   Inches(2.7), Inches(1), Inches(8), Inches(0.6),
                   32, COLORS['text_accent'], bold=False)
    add_title_text(slide, 'Understanding communication from cells to minds',
                   Inches(3.5), Inches(1), Inches(8), Inches(0.4),
                   24, COLORS['text_light'], bold=False)
    
    # Slide 2: Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'We build intelligence that begins where life begins: in cells.',
                   Inches(1), Inches(0.5), Inches(9), Inches(1.2),
                   42, COLORS['text_white'])
    
    body_text = """Healthy cell → immune recognition → cancer immune escape

Human perception → memory → emotional understanding

One unified communication framework."""
    add_body_text(slide, body_text,
                  Inches(2.5), Inches(1.5), Inches(7), Inches(2),
                  26, COLORS['text_light'], PP_ALIGN.CENTER)
    
    # Slide 3: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Communication Failures at Every Scale',
                   Inches(0.5), Inches(1), Inches(8), Inches(0.8),
                   48, COLORS['text_white'])
    
    # Left column
    left_text = """Cancer becomes deadly when it learns to hide from the immune system.

• 600,000+ cancer deaths/year in US alone
• Most immunotherapies fail in solid tumors
• We cannot predict treatment response"""
    add_body_text(slide, left_text,
                  Inches(1.5), Inches(0.5), Inches(4.5), Inches(3),
                  22, COLORS['text_white'])
    
    # Right column
    right_text = """AI becomes untrustworthy when it cannot form meaningful relationships.

• Chatbots have no memory or emotion
• Users don't trust AI systems
• No continuity across interactions"""
    add_body_text(slide, right_text,
                  Inches(1.5), Inches(5.2), Inches(4.5), Inches(3),
                  22, COLORS['text_white'])
    
    # Slide 4: Product 1 - cognisom
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'GPU-Accelerated Multicellular Immune Simulation',
                   Inches(0.5), Inches(1), Inches(8), Inches(0.7),
                   44, COLORS['text_white'])
    
    capabilities = [
        "Intracellular Biochemistry: Transcription, translation, metabolism (2,000-8,000 species/cell)",
        "Immune Recognition: MHC-I presentation, NK/T-cell detection, self vs non-self",
        "Spatial Microenvironment: 3D diffusion, hypoxia gradients, cytokine fields",
        "Cancer Evolution: Clonal mutations, immune evasion, therapy resistance"
    ]
    add_bullet_text(slide, capabilities,
                    Inches(1.4), Inches(0.8), Inches(8.4), Inches(2.5),
                    20, COLORS['text_white'])
    
    add_body_text(slide, 'Initial Focus: Prostate Cancer\nNormal → Oncogenic Stress → Immune Escape → Castration Resistance',
                  Inches(4.2), Inches(2), Inches(6), Inches(0.8),
                  22, COLORS['bio_accent'], PP_ALIGN.CENTER)
    
    # Slide 5: Product 2 - Cogs
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Modular Humanoid Face System with Memory & Emotion',
                   Inches(0.5), Inches(1), Inches(8), Inches(0.7),
                   44, COLORS['text_white'])
    
    cogs_capabilities = [
        "Vision: Face recognition & tracking, familiarity scoring",
        "Hearing: Far-field microphones, sound source localization",
        "Speech & Emotion: Viseme animation, context-aware responses",
        "Memory: pgvector relationship database, Dream Mode nightly learning"
    ]
    add_bullet_text(slide, cogs_capabilities,
                    Inches(1.4), Inches(0.8), Inches(8.4), Inches(2.5),
                    20, COLORS['text_white'])
    
    add_body_text(slide, 'Runs locally on NVIDIA Jetson • Docker microservices (FastAPI) • Real-time perception',
                  Inches(4.2), Inches(2.5), Inches(6), Inches(0.6),
                  20, COLORS['gradient_teal'], PP_ALIGN.CENTER)
    
    # Slide 6: Shared Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'One Communication Model. Two Scales of Life.',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                   44, COLORS['text_white'])
    
    table_data = """Layer          cognisom (Cellular)                    Cogs (Human)
Compute        NVIDIA RTX / H100                    NVIDIA Jetson Orin
Intelligence   Mechanistic biochemical models       Social-emotional reasoning
Memory         pgvector biological states           pgvector relationships
Perception     Reaction-diffusion, signaling        Vision + audio fusion
Adaptation     Clonal evolution, immune evasion     Dream Mode consolidation"""
    
    add_body_text(slide, table_data,
                  Inches(1.5), Inches(0.8), Inches(8), Inches(3.5),
                  18, COLORS['text_white'])
    
    # Slide 7: Why NVIDIA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'NVIDIA Powers Both Platforms',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                   48, COLORS['text_white'])
    
    nvidia_items = [
        "CUDA Kernels: Intracellular SSA, reaction-diffusion PDEs, batched simulation",
        "H100 Scale-Out: Multi-million cell simulations, tumor-immune evolution",
        "Jetson Inference: Real-time humanoid perception, face recognition, audio",
        "TensorRT/Triton: ML surrogate acceleration, 3-10× speedup targets"
    ]
    add_bullet_text(slide, nvidia_items,
                    Inches(1.5), Inches(1), Inches(8), Inches(3.5),
                    22, COLORS['text_white'])
    
    # Slide 8: Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Execution Timeline',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                   48, COLORS['text_white'])
    
    roadmap_text = """Now (0-6 months)
cognisom: Prostate immune surveillance • GPU SSA + spatial diffusion
Cogs: Fully operational prototype • Dream Mode active

Next (6-12 months)
cognisom: H100-scale therapy simulations • Multi-GPU million-cell runs
Cogs: Premium build (AGX Orin) • Enhanced emotion detection

Future (12-24 months)
cognisom: Clinical trial simulation • ML surrogate acceleration
Cogs: Multi-person interaction • Therapeutic applications"""
    
    add_body_text(slide, roadmap_text,
                  Inches(1.3), Inches(1), Inches(8), Inches(3.8),
                  20, COLORS['text_white'])
    
    # Slide 9: Market & Impact
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Transforming Cancer Research & Human-AI Interaction',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.7),
                   42, COLORS['text_white'])
    
    market_text = """Market Opportunity
• Cancer Research: $25B+ annual market, 600k+ deaths/year (US)
• Social Robotics & AI: $35B+ by 2030, elder care, mental health

Our Impact
• Mechanistic cancer understanding, predictive treatment response
• Trustworthy AI companions, accessible assistive technology
• Open-source research tools for the community"""
    
    add_body_text(slide, market_text,
                  Inches(1.4), Inches(1.2), Inches(8), Inches(3.5),
                  22, COLORS['text_white'])
    
    # Slide 10: Competitive Advantage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'What Makes Us Different',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                   48, COLORS['text_white'])
    
    advantage_text = """Existing Approaches → Limitations
• ML-only cancer prediction → No mechanistic interpretability
• Agent-based tumor models → Limited intracellular fidelity
• Standard chatbots/avatars → No memory, no adaptive emotion
• Robotics without cognition → Limited social connection

Our Solution: Mechanistic Biology + Embodied Emotional Intelligence
✓ GPU-accelerated from first principles
✓ Full biochemical fidelity (not phenomenology)
✓ Persistent relational memory • Autonomous learning (Dream Mode)
✓ Open-source, reproducible, standards-based"""
    
    add_body_text(slide, advantage_text,
                  Inches(1.3), Inches(1.2), Inches(8), Inches(3.8),
                  20, COLORS['text_white'])
    
    # Slide 11: Partnership Request
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Join Us in Building the Future',
                   Inches(0.5), Inches(0.5), Inches(9), Inches(0.6),
                   48, COLORS['text_white'])
    
    partnership_text = """What We Need from NVIDIA
✓ H100/A100 compute credits (500-1000 GPU-hours/month)
✓ Jetson optimization guidance (TensorRT, DeepStream)
✓ Technical support (CUDA kernel review, Triton, NCCL)
✓ Co-marketing opportunities (GTC, case studies, blog posts)

What We Offer
✓ Open-source contributions (cognisom platform, CUDA kernels)
✓ Research validation (published results, benchmarks)
✓ Community building (tutorials, collaborations, workshops)
✓ Strategic alignment (showcase NVIDIA's healthcare/biology impact)"""
    
    add_body_text(slide, partnership_text,
                  Inches(1.3), Inches(1.2), Inches(8), Inches(3.8),
                  20, COLORS['text_white'])
    
    # Slide 12: Closing Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_gradient_background(slide, COLORS['primary_dark'], COLORS['gradient_blue'])
    
    add_title_text(slide, 'Understanding Communication from Cells to Minds',
                   Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8),
                   44, COLORS['text_white'])
    
    vision_text = """We are building a future where:

• Cancer treatments are predicted before they're prescribed
• AI systems form genuine, trustworthy relationships
• Biological understanding guides intelligent design
• Open science accelerates discovery for everyone

This is intelligence grounded in life itself."""
    
    add_body_text(slide, vision_text,
                  Inches(1.5), Inches(1.5), Inches(7), Inches(2.5),
                  26, COLORS['text_light'], PP_ALIGN.CENTER)
    
    add_title_text(slide, "Let's build it together on NVIDIA.",
                   Inches(4.2), Inches(2), Inches(6), Inches(0.5),
                   32, COLORS['gradient_teal'])
    
    add_body_text(slide, 'eyentelligence.ai • research@eyentelligence.ai • github.com/eyentelligence',
                  Inches(4.8), Inches(2.5), Inches(6), Inches(0.4),
                  18, COLORS['text_accent'], PP_ALIGN.CENTER)
    
    return prs

if __name__ == '__main__':
    print("Creating eyentelligence pitch deck...")
    prs = create_presentation()
    
    output_path = 'eyentelligence_pitch_deck_B2.pptx'
    prs.save(output_path)
    print(f"✓ Pitch deck saved: {output_path}")
    print(f"✓ Style: Deep Tech Futuristic (B2)")
    print(f"✓ Slides: {len(prs.slides)}")
    print("\nNext steps:")
    print("1. Open in PowerPoint/Keynote/Google Slides")
    print("2. Add logo image to master slides")
    print("3. Add visual elements (neural networks, cellular imagery)")
    print("4. Export to PDF for email distribution")
